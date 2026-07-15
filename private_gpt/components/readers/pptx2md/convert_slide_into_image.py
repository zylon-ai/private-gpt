import logging
import subprocess
import tempfile
from contextlib import suppress
from enum import Enum
from pathlib import Path
from typing import Any

import cv2  # ty:ignore[unresolved-import]
import numpy as np
from PIL import Image, ImageDraw
from pptx import Presentation  # ty:ignore[unresolved-import]
from pptx.enum.shapes import MSO_SHAPE_TYPE  # ty:ignore[unresolved-import]
from pydantic import BaseModel, Field
from scipy import (  # ty:ignore[unresolved-import]
    ndimage,  # type: ignore[import-untyped]
)

from private_gpt.components.ingest.utils import FileInfo
from private_gpt.settings.settings import settings

logger = logging.getLogger(__name__)
logger.setLevel(logging.DEBUG if settings().server.debug_mode else logging.INFO)


class ContentType(Enum):
    TEXT = "text"
    DIAGRAM = "diagram"
    TABLE = "table"
    CHART = "chart"
    IMAGE = "image"
    MEDIA = "media"
    COMPLEX = "complex"


class ContentZone(BaseModel):
    left: int
    top: int
    width: int
    height: int
    content_type: ContentType
    confidence: float = Field(ge=0.0, le=1.0)
    slide_index: int
    zone_id: int


class TextElement(BaseModel):
    left: int
    top: int
    width: int
    height: int
    text: str


class SlideAnalysis(BaseModel):
    slide_index: int
    content_zones: list[ContentZone]
    text_elements: list[TextElement]


class ZoneImageMap(BaseModel):
    slide_index: int
    zone_id: int
    image_path: Path

    class Config:
        arbitrary_types_allowed = True


class ZoneExpansion(BaseModel):
    top: float = Field(ge=0.0, le=100.0)
    bottom: float = Field(ge=0.0, le=100.0)
    left: float = Field(ge=0.0, le=100.0)
    right: float = Field(ge=0.0, le=100.0)


class ExportedImages(BaseModel):
    slide_images: dict[int, Path]
    zone_images: list[ZoneImageMap]
    content_zones: list[ContentZone]
    export_directory: Path

    class Config:
        arbitrary_types_allowed = True


class TransformConfig(BaseModel):
    confidence_threshold: float = Field(default=0.6, ge=0.0, le=1.0)
    libreoffice_timeout: int = Field(default=60, gt=0)
    min_zone_size_ratio: float = Field(default=0.05, ge=0.0, le=1.0)
    grid_width: int = Field(default=120, gt=0)
    grid_height: int = Field(default=80, gt=0)
    debug_mode: bool = Field(default=False)
    zones_enabled: bool = Field(default=True)
    zone_expansion_percentage: ZoneExpansion = Field(
        default=ZoneExpansion(top=20.0, bottom=20.0, left=15.0, right=15.0)
    )
    overlap_threshold_for_duplicates: float = Field(default=0.9, ge=0.0, le=1.0)
    shape_removal_overlap_threshold: float = Field(default=0.8, ge=0.0, le=1.0)


class EnhancedPPTXTransform:
    def __init__(self, config: TransformConfig | None = None) -> None:
        self.config = config or TransformConfig()
        self.presentation: Any | None = None

    async def transform_file(
        self, file_info: FileInfo
    ) -> tuple[FileInfo, ExportedImages | None]:
        original_path = file_info.file_data
        self.presentation = Presentation(str(original_path))

        export_temp_dir = Path(tempfile.mkdtemp(prefix="pptx_export_"))
        logger.debug(f"Exporting images to: {export_temp_dir}")

        with tempfile.TemporaryDirectory() as temp_dir:
            temp_path = Path(temp_dir)

            slide_images = self._convert_pptx_to_pngs(original_path, temp_path)
            if not slide_images:
                logger.warning("No slide images generated, skipping conversion")
                return file_info, None

            if not self.config.zones_enabled:
                exported = self._export_images(slide_images, [], [], export_temp_dir)
                return file_info, exported

            content_zones = await self._detect_content_zones(
                self.presentation, slide_images
            )
            if not content_zones:
                logger.debug(
                    "No content zones detected for conversion, returning original"
                )
                return file_info, None

            zone_images = self._extract_zone_images(
                slide_images, content_zones, temp_path, self.presentation
            )

            grouped_zone_images_by_slide = await self._group_zone_images_by_slide(
                zone_images, content_zones
            )

            zone_images = [
                zi
                for slide_zis in grouped_zone_images_by_slide.values()
                for zi in slide_zis
            ]

            exported = self._export_images(
                slide_images, zone_images, content_zones, export_temp_dir
            )

            modified_path = self._create_modified_pptx_with_zones(
                original_path, content_zones, zone_images
            )

            file_info_copy = file_info.model_copy(deep=True)
            file_info_copy.file_data = modified_path
            return file_info_copy, exported

    def _export_images(
        self,
        slide_images: dict[int, Path],
        zone_images: list[ZoneImageMap],
        content_zones: list[ContentZone],
        export_dir: Path,
    ) -> ExportedImages:
        """Export slide and zone images to the specified directory."""
        import shutil

        # Create subdirectories
        slides_dir = export_dir / "slides"
        zones_dir = export_dir / "zones"
        slides_dir.mkdir(parents=True, exist_ok=True)
        zones_dir.mkdir(parents=True, exist_ok=True)

        # Export slide images
        exported_slide_images: dict[int, Path] = {}
        for slide_idx, slide_path in slide_images.items():
            dest_path = slides_dir / f"slide_{slide_idx:02d}.png"
            shutil.copy2(slide_path, dest_path)
            exported_slide_images[slide_idx] = dest_path
            logger.debug(f"Exported slide {slide_idx} to {dest_path}")

        # Export zone images
        exported_zone_images: list[ZoneImageMap] = []
        for zone_image in zone_images:
            dest_path = zones_dir / zone_image.image_path.name
            shutil.copy2(zone_image.image_path, dest_path)
            exported_zone_images.append(
                ZoneImageMap(
                    slide_index=zone_image.slide_index,
                    zone_id=zone_image.zone_id,
                    image_path=dest_path,
                )
            )
            logger.debug(f"Exported zone image to {dest_path}")

        logger.debug(
            f"Exported {len(exported_slide_images)} slide images and "
            f"{len(exported_zone_images)} zone images to {export_dir}"
        )

        return ExportedImages(
            slide_images=exported_slide_images,
            zone_images=exported_zone_images,
            content_zones=content_zones,
            export_directory=export_dir,
        )

    async def _detect_content_zones(
        self, prs: Any, slide_images: dict[int, Path]
    ) -> list[ContentZone]:
        all_zones: list[ContentZone] = []

        for slide_idx in range(len(prs.slides)):
            if slide_idx not in slide_images:
                continue

            slide = prs.slides[slide_idx]
            text_elements = self._extract_text_elements(slide)

            zones = self._detect_non_text_zones(
                text_elements, prs.slide_width, prs.slide_height, slide_idx
            )

            complex_zones = self._detect_existing_complex_shapes(slide, slide_idx)
            zones.extend(complex_zones)

            filtered_zones = await self._filter_zones_advanced(
                zones, prs.slide_width, prs.slide_height
            )
            all_zones.extend(filtered_zones)

        logger.debug(
            f"Detected {len(all_zones)} content zones across {len(slide_images)} slides"
        )
        return all_zones

    def _extract_text_elements(self, slide: Any) -> list[TextElement]:
        text_elements: list[TextElement] = []

        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text.strip():
                text_elements.append(
                    TextElement(
                        left=shape.left,
                        top=shape.top,
                        width=shape.width,
                        height=shape.height,
                        text=shape.text.strip(),
                    )
                )

        return text_elements

    def _detect_non_text_zones(
        self,
        text_elements: list[TextElement],
        slide_width: int,
        slide_height: int,
        slide_index: int,
    ) -> list[ContentZone]:

        occupancy_map = np.zeros(
            (self.config.grid_height, self.config.grid_width), dtype=np.float64
        )

        for element in text_elements:
            x_start = max(0, int((element.left / slide_width) * self.config.grid_width))
            x_end = min(
                self.config.grid_width,
                int(
                    ((element.left + element.width) / slide_width)
                    * self.config.grid_width
                ),
            )
            y_start = max(
                0, int((element.top / slide_height) * self.config.grid_height)
            )
            y_end = min(
                self.config.grid_height,
                int(
                    ((element.top + element.height) / slide_height)
                    * self.config.grid_height
                ),
            )

            occupancy_map[y_start:y_end, x_start:x_end] = 1

        occupancy_map = ndimage.gaussian_filter(occupancy_map, sigma=1.0)
        empty_map_float = (1 - occupancy_map) > 0.5
        empty_map = empty_map_float.astype(np.uint8)

        kernel = np.ones((3, 3), np.uint8)
        empty_map = cv2.morphologyEx(empty_map, cv2.MORPH_CLOSE, kernel)  # type: ignore[assignment]
        empty_map = cv2.morphologyEx(empty_map, cv2.MORPH_OPEN, kernel)  # type: ignore[assignment]

        num_labels, labels = cv2.connectedComponents(empty_map)

        zones: list[ContentZone] = []
        zone_id = 0

        for label in range(1, num_labels):
            mask = labels == label
            coords_y, coords_x = np.where(mask)

            if len(coords_y) == 0:
                continue

            min_y, max_y = int(coords_y.min()), int(coords_y.max())
            min_x, max_x = int(coords_x.min()), int(coords_x.max())

            left = int((min_x / self.config.grid_width) * slide_width)
            top = int((min_y / self.config.grid_height) * slide_height)
            width = int(((max_x - min_x) / self.config.grid_width) * slide_width)
            height = int(((max_y - min_y) / self.config.grid_height) * slide_height)

            content_type, confidence = self._classify_content_zone(
                width, height, slide_width, slide_height
            )

            if confidence > 0.3:
                zone = ContentZone(
                    left=left,
                    top=top,
                    width=width,
                    height=height,
                    content_type=content_type,
                    confidence=confidence,
                    slide_index=slide_index,
                    zone_id=zone_id,
                )
                zones.append(zone)
                zone_id += 1

        return zones

    def _classify_content_zone(
        self, width: int, height: int, slide_width: int, slide_height: int
    ) -> tuple[ContentType, float]:

        aspect_ratio = width / height
        area_ratio = (width * height) / (slide_width * slide_height)

        if area_ratio > 0.25:
            if 0.8 <= aspect_ratio <= 1.2:
                return ContentType.DIAGRAM, 0.8
            elif aspect_ratio > 2.5:
                return ContentType.CHART, 0.85
            else:
                return ContentType.COMPLEX, 0.7

        elif area_ratio > 0.08:
            if aspect_ratio > 3:
                return ContentType.TABLE, 0.9
            elif 0.5 <= aspect_ratio <= 2:
                return ContentType.DIAGRAM, 0.75
            else:
                return ContentType.COMPLEX, 0.6

        else:
            if aspect_ratio > 2.5:
                return ContentType.TABLE, 0.7
            else:
                return ContentType.DIAGRAM, 0.5

    def _detect_existing_complex_shapes(
        self, slide: Any, slide_index: int
    ) -> list[ContentZone]:
        zones: list[ContentZone] = []
        zone_id = 1000

        safe_shape_types: list[Any] = []
        type_mapping: dict[Any, ContentType] = {}

        try:
            safe_shape_types.append(MSO_SHAPE_TYPE.CHART)
            type_mapping[MSO_SHAPE_TYPE.CHART] = ContentType.CHART
        except AttributeError:
            pass

        try:
            safe_shape_types.append(MSO_SHAPE_TYPE.PICTURE)
            type_mapping[MSO_SHAPE_TYPE.PICTURE] = ContentType.IMAGE
        except AttributeError:
            pass

        try:
            safe_shape_types.append(MSO_SHAPE_TYPE.MEDIA)
            type_mapping[MSO_SHAPE_TYPE.MEDIA] = ContentType.MEDIA
        except AttributeError:
            pass

        try:
            safe_shape_types.append(MSO_SHAPE_TYPE.TABLE)
            type_mapping[MSO_SHAPE_TYPE.TABLE] = ContentType.TABLE
        except AttributeError:
            pass

        for smart_art_name in ["SMART_ART", "SMARTART", "SMART_ART_GRAPHIC"]:
            try:
                smart_art_type = getattr(MSO_SHAPE_TYPE, smart_art_name)
                safe_shape_types.append(smart_art_type)
                type_mapping[smart_art_type] = ContentType.DIAGRAM
                break
            except AttributeError:
                continue

        for shape in slide.shapes:
            try:
                if (
                    hasattr(shape, "shape_type")
                    and shape.shape_type in safe_shape_types
                ):
                    content_type = type_mapping.get(
                        shape.shape_type, ContentType.COMPLEX
                    )

                    if not all(
                        hasattr(shape, attr)
                        for attr in ["left", "top", "width", "height"]
                    ):
                        continue

                    zone = ContentZone(
                        left=shape.left,
                        top=shape.top,
                        width=shape.width,
                        height=shape.height,
                        content_type=content_type,
                        confidence=0.95,
                        slide_index=slide_index,
                        zone_id=zone_id,
                    )
                    zones.append(zone)
                    zone_id += 1

                elif hasattr(shape, "shape_type"):
                    shape_type_str = str(shape.shape_type)

                    if any(
                        keyword in shape_type_str.upper()
                        for keyword in [
                            "CHART",
                            "PICTURE",
                            "IMAGE",
                            "MEDIA",
                            "SMART",
                            "TABLE",
                            "DIAGRAM",
                        ]
                    ):
                        if all(
                            hasattr(shape, attr)
                            for attr in ["left", "top", "width", "height"]
                        ):
                            zone = ContentZone(
                                left=shape.left,
                                top=shape.top,
                                width=shape.width,
                                height=shape.height,
                                content_type=ContentType.COMPLEX,
                                confidence=0.8,
                                slide_index=slide_index,
                                zone_id=zone_id,
                            )
                            zones.append(zone)
                            zone_id += 1

            except Exception as e:
                logger.warning(f"Error processing shape on slide {slide_index}: {e}")
                continue

        return zones

    def _filter_zones(
        self, zones: list[ContentZone], slide_width: int, slide_height: int
    ) -> list[ContentZone]:
        filtered: list[ContentZone] = []
        slide_area = slide_width * slide_height
        min_area = slide_area * self.config.min_zone_size_ratio

        for zone in zones:
            zone_area = zone.width * zone.height

            if (
                zone.confidence >= self.config.confidence_threshold
                and zone_area >= min_area
            ):
                filtered.append(zone)

        logger.debug(f"Filtered {len(filtered)} zones from {len(zones)} candidates")
        return filtered

    def _extract_zone_images(
        self,
        slide_images: dict[int, Path],
        content_zones: list[ContentZone],
        temp_dir: Path,
        presentation: Any,
    ) -> list[ZoneImageMap]:

        zone_images: list[ZoneImageMap] = []
        actual_slide_width = presentation.slide_width
        actual_slide_height = presentation.slide_height

        for zone in content_zones:
            if zone.slide_index not in slide_images:
                continue

            try:
                slide_image_path = slide_images[zone.slide_index]
                slide_img = Image.open(slide_image_path)
                img_width, img_height = slide_img.size

                if self.config.debug_mode:
                    logger.debug(
                        f"Processing zone {zone.zone_id} on slide {zone.slide_index}"
                    )
                    logger.debug(
                        f"Zone bounds (EMU): left={zone.left}, top={zone.top}, width={zone.width}, height={zone.height}"
                    )
                    logger.debug(
                        f"Slide size (EMU): {actual_slide_width} x {actual_slide_height}"
                    )
                    logger.debug(f"Image size (px): {img_width} x {img_height}")

                left_rel = max(0.0, min(1.0, zone.left / actual_slide_width))
                top_rel = max(0.0, min(1.0, zone.top / actual_slide_height))
                right_rel = max(
                    0.0, min(1.0, (zone.left + zone.width) / actual_slide_width)
                )
                bottom_rel = max(
                    0.0, min(1.0, (zone.top + zone.height) / actual_slide_height)
                )

                left_px = int(left_rel * img_width)
                top_px = int(top_rel * img_height)
                right_px = int(right_rel * img_width)
                bottom_px = int(bottom_rel * img_height)

                left_px = max(0, min(left_px, img_width - 1))
                top_px = max(0, min(top_px, img_height - 1))
                right_px = max(left_px + 1, min(right_px, img_width))
                bottom_px = max(top_px + 1, min(bottom_px, img_height))

                if (right_px - left_px) < 10 or (bottom_px - top_px) < 10:
                    logger.warning(
                        f"Zone {zone.zone_id} too small, expanding to minimum size"
                    )
                    center_x = (left_px + right_px) // 2
                    center_y = (top_px + bottom_px) // 2
                    left_px = max(0, center_x - 25)
                    right_px = min(img_width, center_x + 25)
                    top_px = max(0, center_y - 25)
                    bottom_px = min(img_height, center_y + 25)

                crop_box = (left_px, top_px, right_px, bottom_px)
                cropped = slide_img.crop(crop_box)

                zone_image_path = (
                    temp_dir
                    / f"zone_s{zone.slide_index}_z{zone.zone_id}_{zone.content_type.value}.png"
                )
                cropped.save(zone_image_path)

                zone_images.append(
                    ZoneImageMap(
                        slide_index=zone.slide_index,
                        zone_id=zone.zone_id,
                        image_path=zone_image_path,
                    )
                )

                logger.debug(
                    f"Created: {zone_image_path.name} ({cropped.size[0]}x{cropped.size[1]} px)"
                )

                if self.config.debug_mode:
                    debug_path = (
                        temp_dir
                        / f"debug_slide_{zone.slide_index}_zone_{zone.zone_id}.png"
                    )
                    debug_img = slide_img.copy()
                    draw = ImageDraw.Draw(debug_img)
                    draw.rectangle(
                        [left_px, top_px, right_px, bottom_px], outline="red", width=3
                    )
                    debug_img.save(debug_path)

            except Exception as e:
                logger.error(
                    f"Error extracting zone {zone.zone_id} from slide {zone.slide_index}: {e}"
                )
                continue

        logger.debug(f"Successfully extracted {len(zone_images)} zone images")
        return zone_images

    def _create_modified_pptx_with_zones(
        self,
        original_path: Path,
        content_zones: list[ContentZone],
        zone_images: list[ZoneImageMap],
    ) -> Path:
        prs = Presentation(str(original_path))

        zone_image_lookup: dict[tuple[int, int], Path] = {
            (zi.slide_index, zi.zone_id): zi.image_path for zi in zone_images
        }

        zones_by_slide: dict[int, list[ContentZone]] = {}
        for zone in content_zones:
            if zone.slide_index not in zones_by_slide:
                zones_by_slide[zone.slide_index] = []
            zones_by_slide[zone.slide_index].append(zone)

        for slide_idx, zones in zones_by_slide.items():
            if slide_idx >= len(prs.slides):
                continue

            slide = prs.slides[slide_idx]
            self._remove_overlapping_shapes(slide, zones)

            for zone in zones:
                zone_key = (zone.slide_index, zone.zone_id)
                if zone_key in zone_image_lookup:
                    self._remove_original_content(slide, zone)
                    self._add_zone_image(slide, zone, zone_image_lookup[zone_key])

        modified_path = (
            original_path.parent
            / f"{original_path.stem}_zone_converted{original_path.suffix}"
        )
        prs.save(str(modified_path))
        logger.debug(f"Modified presentation saved to: {modified_path}")
        return modified_path

    def _remove_overlapping_shapes(self, slide: Any, zones: list[ContentZone]) -> None:
        shapes_to_remove: list[Any] = []

        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text.strip():
                continue

            for zone in zones:
                if self._shapes_overlap(shape, zone, threshold=0.7):
                    shapes_to_remove.append(shape)
                    break

        for shape in shapes_to_remove:
            with suppress(Exception):
                sp = shape._element
                sp.getparent().remove(sp)

    def _shapes_overlap(
        self, shape: Any, zone: ContentZone, threshold: float = 0.5
    ) -> bool:
        if not all(hasattr(shape, attr) for attr in ["left", "top", "width", "height"]):
            return False

        left1, top1 = shape.left, shape.top
        right1, bottom1 = shape.left + shape.width, shape.top + shape.height

        left2, top2 = zone.left, zone.top
        right2, bottom2 = zone.left + zone.width, zone.top + zone.height

        left_int = max(left1, left2)
        top_int = max(top1, top2)
        right_int = min(right1, right2)
        bottom_int = min(bottom1, bottom2)

        if left_int >= right_int or top_int >= bottom_int:
            return False

        intersection_area = (right_int - left_int) * (bottom_int - top_int)
        shape_area = shape.width * shape.height

        if shape_area == 0:
            return False

        overlap_ratio = intersection_area / shape_area
        return bool(overlap_ratio >= threshold)

    def _add_zone_image(self, slide: Any, zone: ContentZone, image_path: Path) -> None:
        try:
            slide.shapes.add_picture(
                str(image_path), zone.left, zone.top, zone.width, zone.height
            )
        except Exception as e:
            logger.warning(f"Could not add zone image: {e}")

    def _convert_pptx_to_pngs(self, pptx_path: Path, temp_dir: Path) -> dict[int, Path]:
        pdf_path = temp_dir / f"{pptx_path.stem}.pdf"

        cmd_pdf = [
            "soffice",
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(temp_dir),
            str(pptx_path),
        ]

        try:
            result = subprocess.run(
                cmd_pdf, capture_output=True, timeout=self.config.libreoffice_timeout
            )
            if result.returncode != 0 or not pdf_path.exists():
                stderr_text = (
                    result.stderr.decode("utf-8", errors="replace")
                    if result.stderr
                    else "No error output"
                )
                logger.error(f"Failed to convert PPTX to PDF: {stderr_text}")
                return {}
        except subprocess.TimeoutExpired:
            logger.error("PPTX to PDF conversion timed out")
            return {}

        output_pattern = str(temp_dir / "slide-%02d.png")
        cmd_png = [
            "gs",
            "-sDEVICE=pngalpha",
            f"-o{output_pattern}",
            "-r144",
            str(pdf_path),
        ]

        try:
            result = subprocess.run(
                cmd_png, capture_output=True, timeout=self.config.libreoffice_timeout
            )
            if result.returncode != 0:
                stderr_text = (
                    result.stderr.decode("utf-8", errors="replace")
                    if result.stderr
                    else "No error output"
                )
                logger.error(f"Failed to convert PDF to PNG: {stderr_text}")
                return {}
        except subprocess.TimeoutExpired:
            logger.error("PDF to PNG conversion timed out")
            return {}

        png_files = sorted(temp_dir.glob("slide-*.png"))
        slide_images: dict[int, Path] = dict(enumerate(png_files))
        logger.debug(f"Generated {len(slide_images)} slide images")
        return slide_images

    async def _group_zone_images_by_slide(
        self, zone_images: list[ZoneImageMap], content_zones: list[ContentZone]
    ) -> dict[int, list[ZoneImageMap]]:
        """Extract and group zone images by slide index with filtering logic."""
        grouped_zone_images_by_slide: dict[int, list[ZoneImageMap]] = {}

        # Group zone images by slide
        for zi in zone_images:
            if zi.slide_index not in grouped_zone_images_by_slide:
                grouped_zone_images_by_slide[zi.slide_index] = []
            grouped_zone_images_by_slide[zi.slide_index].append(zi)

        # Apply filtering logic for each slide
        for slide_idx, zis in grouped_zone_images_by_slide.items():
            logger.debug(f"Slide {slide_idx} has {len(zis)} zone images extracted.")

            has_non_general_zone = any(
                z.zone_id != 0 for z in content_zones if z.slide_index == slide_idx
            )

            if has_non_general_zone:
                logger.debug(
                    f"Slide {slide_idx} has non-general content zones detected."
                )
                # Remove general zones (zone_id == 0)
                grouped_zone_images_by_slide[slide_idx] = [
                    z for z in zis if z.zone_id != 0
                ]
                logger.debug(
                    f"After filtering, slide {slide_idx} has "
                    f"{len(grouped_zone_images_by_slide[slide_idx])} zone images."
                )

        return grouped_zone_images_by_slide

    def _expand_zone(
        self, zone: ContentZone, slide_width: int, slide_height: int
    ) -> ContentZone:
        """Expand a zone by the configured percentage in all directions."""
        top_expansion_factor = self.config.zone_expansion_percentage.top / 100.0
        bottom_expansion_factor = self.config.zone_expansion_percentage.bottom / 100.0
        left_expansion_factor = self.config.zone_expansion_percentage.left / 100.0
        right_expansion_factor = self.config.zone_expansion_percentage.right / 100.0

        # Calculate expansion amounts
        top_expansion = int(zone.height * top_expansion_factor)
        bottom_expansion = int(zone.height * bottom_expansion_factor)
        left_expansion = int(zone.width * left_expansion_factor)
        right_expansion = int(zone.width * right_expansion_factor)

        # Calculate new bounds
        new_top = max(0, zone.top - top_expansion)
        new_bottom = max(0, zone.top + zone.height + bottom_expansion)
        new_left = max(0, zone.left - left_expansion)
        new_right = max(0, zone.left + zone.width + right_expansion)

        if new_top < 0:
            new_top = 0
        if new_left < 0:
            new_left = 0
        if new_bottom > slide_height:
            new_bottom = slide_height
        if new_right > slide_width:
            new_right = slide_width

        new_width = min(slide_width - new_left, new_right - new_left)
        new_height = min(slide_height - new_top, new_bottom - new_top)

        return ContentZone(
            left=new_left,
            top=new_top,
            width=new_width,
            height=new_height,
            content_type=zone.content_type,
            confidence=zone.confidence,
            slide_index=zone.slide_index,
            zone_id=zone.zone_id,
        )

    def _calculate_zone_overlap(self, zone1: ContentZone, zone2: ContentZone) -> float:
        """Calculate the overlap ratio between two zones."""
        if zone1.slide_index != zone2.slide_index:
            return 0.0

        # Calculate intersection boundaries
        left_int = max(zone1.left, zone2.left)
        top_int = max(zone1.top, zone2.top)
        right_int = min(zone1.left + zone1.width, zone2.left + zone2.width)
        bottom_int = min(zone1.top + zone1.height, zone2.top + zone2.height)

        # No intersection if boundaries don't overlap
        if left_int >= right_int or top_int >= bottom_int:
            return 0.0

        # Calculate areas
        intersection_area = (right_int - left_int) * (bottom_int - top_int)
        zone1_area = zone1.width * zone1.height
        zone2_area = zone2.width * zone2.height

        if zone1_area == 0 or zone2_area == 0:
            return 0.0

        # Return the overlap ratio relative to the smaller zone
        smaller_area = min(zone1_area, zone2_area)
        return intersection_area / smaller_area

    async def _filter_zones_advanced(
        self, zones: list[ContentZone], slide_width: int, slide_height: int
    ) -> list[ContentZone]:
        # Step 1: Apply basic filtering
        basic_filtered = self._filter_zones(zones, slide_width, slide_height)

        # Step 2: Expand zones
        expanded_zones = [
            self._expand_zone(zone, slide_width, slide_height)
            for zone in basic_filtered
        ]

        # Step 3: Group zones by slide for processing
        zones_by_slide: dict[int, list[ContentZone]] = {}
        for zone in expanded_zones:
            if zone.slide_index not in zones_by_slide:
                zones_by_slide[zone.slide_index] = []
            zones_by_slide[zone.slide_index].append(zone)

        final_filtered_zones: list[ContentZone] = []

        # Step 4: Process each slide separately
        for slide_idx, slide_zones in zones_by_slide.items():
            # Check if there are non-general zones (zone_id != 0)
            non_general_zones = [z for z in slide_zones if z.zone_id != 0]
            general_zones = [z for z in slide_zones if z.zone_id == 0]

            # If we have non-general zones, filter out general zones
            if non_general_zones:
                logger.debug(
                    f"Slide {slide_idx}: Found {len(non_general_zones)} non-general zones, "
                    f"removing {len(general_zones)} general zones"
                )
                working_zones = non_general_zones
            else:
                working_zones = slide_zones

            # Step 5: Remove high-overlap duplicates
            filtered_slide_zones = self._remove_duplicate_zones(working_zones)
            final_filtered_zones.extend(filtered_slide_zones)

        logger.debug(
            f"Filtering: {len(zones)} -> {len(basic_filtered)} -> "
            f"{len(expanded_zones)} -> {len(final_filtered_zones)} zones"
        )

        return final_filtered_zones

    def _remove_duplicate_zones(self, zones: list[ContentZone]) -> list[ContentZone]:
        """Remove zones with high overlap, keeping the one with higher confidence."""
        if len(zones) <= 1:
            return zones

        # Sort by confidence (descending) to prioritize higher confidence zones
        sorted_zones = sorted(zones, key=lambda z: z.confidence, reverse=True)
        filtered_zones: list[ContentZone] = []

        for current_zone in sorted_zones:
            is_duplicate = False

            for existing_zone in filtered_zones:
                overlap_ratio = self._calculate_zone_overlap(
                    current_zone, existing_zone
                )

                if overlap_ratio >= self.config.overlap_threshold_for_duplicates:
                    logger.debug(
                        f"Removing duplicate zone {current_zone.zone_id} "
                        f"(overlap: {overlap_ratio:.2f} with zone {existing_zone.zone_id})"
                    )
                    is_duplicate = True
                    break

            if not is_duplicate:
                filtered_zones.append(current_zone)

        logger.debug(
            f"Duplicate removal: {len(zones)} -> {len(filtered_zones)} zones "
            f"(removed {len(zones) - len(filtered_zones)} duplicates)"
        )

        return filtered_zones

    def _remove_original_content(self, slide: Any, zone: ContentZone) -> int:
        """Remove original content that overlaps with the zone to avoid duplication."""
        shapes_to_remove = []

        for shape in slide.shapes:
            try:
                overlap = self._shapes_overlap(
                    shape, zone, threshold=self.config.shape_removal_overlap_threshold
                )
                if overlap:
                    shapes_to_remove.append(shape)

            except Exception:
                continue

        # Remove the shapes
        for shape in shapes_to_remove:
            try:
                sp = shape._element
                sp.getparent().remove(sp)
            except Exception:
                continue

        logger.debug(
            f"Removed {len(shapes_to_remove)} overlapping shapes for zone {zone.zone_id} on slide {zone.slide_index}"
        )
        return len(shapes_to_remove)


class PPTXSlideToImageTransform(EnhancedPPTXTransform):
    """Compatibility wrapper that maintains the original interface."""

    def __init__(
        self, confidence_threshold: float = 0.7, libreoffice_timeout: int = 60
    ) -> None:
        config = TransformConfig(
            confidence_threshold=confidence_threshold,
            libreoffice_timeout=libreoffice_timeout,
            zones_enabled=False,
        )
        super().__init__(config)


class PPTXSlideToImageDeepTransform(EnhancedPPTXTransform):
    """Compatibility wrapper that maintains the original interface."""

    def __init__(
        self, confidence_threshold: float = 0.7, libreoffice_timeout: int = 60
    ) -> None:
        config = TransformConfig(
            confidence_threshold=confidence_threshold,
            libreoffice_timeout=libreoffice_timeout,
            zones_enabled=True,
        )
        super().__init__(config)
